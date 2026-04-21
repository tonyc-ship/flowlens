#!/usr/bin/env python3
"""
Extract text content from office documents (.docx, .pptx, .pdf)
and save as structured Markdown.

Usage:
    python extract_document.py <input_file> <output_md_path>

Supported formats:
    - .docx (Word)
    - .pptx (PowerPoint)
    - .pdf  (PDF)
"""

import sys
import os
from pathlib import Path
from datetime import datetime


def extract_docx(filepath: str) -> str:
    from docx import Document
    doc = Document(filepath)

    sections = []
    current_heading = None

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        if para.style.name.startswith("Heading"):
            level = 1
            try:
                level = int(para.style.name.replace("Heading ", "").replace("Heading", "1"))
            except ValueError:
                level = 1
            sections.append(f"\n{'#' * (level + 1)} {text}\n")
        elif para.style.name == "List Paragraph" or para.style.name.startswith("List"):
            sections.append(f"- {text}")
        else:
            sections.append(text)

    for table in doc.tables:
        rows = []
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
        sections.append("\n" + "\n".join(rows) + "\n")

    return "\n".join(sections)


def extract_pptx(filepath: str) -> str:
    from pptx import Presentation
    prs = Presentation(filepath)

    sections = []

    for i, slide in enumerate(prs.slides, 1):
        slide_title = None
        slide_content = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if shape.shape_type == 13 or (hasattr(shape, "is_placeholder") and shape.placeholder_format and shape.placeholder_format.idx == 0):
                        slide_title = text
                    else:
                        slide_content.append(text)
            if shape.has_table:
                rows = []
                for ri, row in enumerate(shape.table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                    if ri == 0:
                        rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                slide_content.append("\n".join(rows))

        title_text = slide_title or f"Slide {i}"
        sections.append(f"\n## {title_text}\n")
        for line in slide_content:
            if line != title_text:
                sections.append(line)

    return "\n".join(sections)


def extract_pdf(filepath: str) -> str:
    import pdfplumber

    sections = []

    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text:
                sections.append(f"\n## Page {i}\n")
                sections.append(text.strip())

            tables = page.extract_tables()
            for table in tables:
                if not table:
                    continue
                rows = []
                for ri, row in enumerate(table):
                    cells = [str(cell or "").strip() for cell in row]
                    rows.append("| " + " | ".join(cells) + " |")
                    if ri == 0:
                        rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                sections.append("\n" + "\n".join(rows) + "\n")

    return "\n".join(sections)


EXTRACTORS = {
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".pdf": extract_pdf,
}


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract_document.py <input_file> [output_md_path]")
        print("Supported: .docx, .pptx, .pdf")
        sys.exit(1)

    input_file = sys.argv[1]
    if not os.path.exists(input_file):
        print(f"Error: File not found: {input_file}")
        sys.exit(1)

    ext = Path(input_file).suffix.lower()
    if ext not in EXTRACTORS:
        print(f"Error: Unsupported format '{ext}'. Supported: {', '.join(EXTRACTORS.keys())}")
        sys.exit(1)

    if len(sys.argv) >= 3:
        output_path = sys.argv[2]
    else:
        stem = Path(input_file).stem
        date_prefix = datetime.now().strftime("%y%m%d")
        output_path = f"[{date_prefix}]{stem}.md"

    print(f"Extracting {ext} → {output_path}")
    content = EXTRACTORS[ext](input_file)

    filename = Path(input_file).name
    header = f"---\ntitle: \"{Path(input_file).stem}\"\nsource_file: \"{filename}\"\nformat: \"{ext[1:]}\"\nextracted: \"{datetime.now().strftime('%Y-%m-%d %H:%M')}\"\n---\n\n# {Path(input_file).stem}\n"

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(header + "\n" + content)

    print(f"Done. Saved to: {output_path}")
    print(f"Content length: {len(content)} chars")


if __name__ == "__main__":
    main()
